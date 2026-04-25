"""Generate VibeMatch AI presentation as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
PURPLE      = RGBColor(0x6C, 0x47, 0xFF)   # brand accent
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)   # slide background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill=None, line=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, left, top, width, height,
                   size=20, color=WHITE, bullet="•"):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(8)


def pill(slide, text, left, top, width=2.2, height=0.45,
         fill=PURPLE, tsize=16, tcolor=WHITE):
    s = box(slide, left, top, width, height, fill=fill)
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(tsize)
    run.font.bold = True
    run.font.color.rgb = tcolor
    return s


def accent_bar(slide, top=0.55):
    """Thin purple rule under title area."""
    bar = slide.shapes.add_shape(1,
        Inches(0.5), Inches(top), Inches(12.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1)

# big music note watermark
txt(s1, "🎵", 9.5, 1.5, 3, 4, size=180, color=RGBColor(0x2A, 0x2A, 0x4E),
    align=PP_ALIGN.CENTER)

txt(s1, "VibeMatch AI", 0.6, 1.8, 9, 1.4, size=60, bold=True,
    color=WHITE, align=PP_ALIGN.LEFT)

txt(s1, "A Music Recommendation System Powered by RAG + Claude",
    0.6, 3.1, 10, 0.7, size=24, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

accent_bar(s1, top=4.1)

txt(s1, "CodePath AI110  ·  Applied AI Systems Project",
    0.6, 4.3, 9, 0.5, size=18, color=LIGHT_GRAY, italic=True)

pill(s1, "github.com/yarinacs/applied-ai-system-project",
     0.6, 6.6, 6.8, 0.45, fill=RGBColor(0x2A, 0x2A, 0x4E), tsize=14,
     tcolor=ACCENT_TEAL)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem / What I Built
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2)

txt(s2, "The Problem With Old-Style Recommenders",
    0.5, 0.25, 12, 0.65, size=32, bold=True, color=WHITE)
accent_bar(s2, top=0.95)

# Left panel
box(s2, 0.5, 1.15, 5.8, 5.5, fill=RGBColor(0x25, 0x25, 0x40), line=PURPLE)
txt(s2, "Old approach (Module 3)",
    0.7, 1.25, 5.4, 0.5, size=18, bold=True, color=PURPLE)
add_bullet_box(s2,
    ["Fill in a form: genre, mood, energy",
     "System matches on exact labels",
     "No understanding of what you mean",
     "Same inputs → same outputs, always"],
    0.7, 1.85, 5.4, 3.0, size=19, color=LIGHT_GRAY)

# Right panel
box(s2, 6.9, 1.15, 5.9, 5.5, fill=RGBColor(0x1E, 0x2A, 0x3A), line=ACCENT_TEAL)
txt(s2, "VibeMatch AI",
    7.1, 1.25, 5.5, 0.5, size=18, bold=True, color=ACCENT_TEAL)
add_bullet_box(s2,
    ['Type: "chill for late-night studying"',
     "System understands the intent",
     "RAG retrieves semantically similar songs",
     "Claude explains why it picked that song"],
    7.1, 1.85, 5.5, 3.0, size=19, color=WHITE)

txt(s2, "→", 6.1, 3.5, 0.8, 0.8, size=42, bold=True,
    color=PURPLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3)

txt(s3, "System Architecture", 0.5, 0.25, 12, 0.65,
    size=32, bold=True, color=WHITE)
accent_bar(s3, top=0.95)

# Insert diagram image
diagram_path = os.path.join(os.path.dirname(__file__), "mermaid_diagram.jpg")
if os.path.exists(diagram_path):
    s3.shapes.add_picture(diagram_path,
        Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.6))

# Step pills on the right
steps = [
    ("1", "Query → vector embedding",     PURPLE),
    ("2", "RAG retrieves top 5 songs",    RGBColor(0x00, 0xAA, 0xCC)),
    ("3", "Claude picks best match",      ACCENT_TEAL),
    ("4", "Confidence score shown",       RGBColor(0xFF, 0xA5, 0x00)),
]
for i, (num, label, color) in enumerate(steps):
    top = 1.3 + i * 1.3
    pill(s3, num, 9.0, top, 0.45, 0.45, fill=color, tsize=18)
    txt(s3, label, 9.6, top, 3.5, 0.6, size=17, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Live Demo
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4)

txt(s4, "Live Demo", 0.5, 0.25, 12, 0.65, size=32, bold=True, color=WHITE)
accent_bar(s4, top=0.95)

txt(s4, "3 queries — watch how behavior changes",
    0.5, 1.1, 12, 0.5, size=20, color=LIGHT_GRAY, italic=True)

demos = [
    ('"something chill and focused for late night studying"',
     "Focus Flow  ·  High confidence  ·  YouTube autoplays",
     ACCENT_TEAL),
    ('"high energy upbeat song for a morning run"',
     "Pulse Horizon  ·  High confidence  ·  0.95 energy, 140 BPM",
     PURPLE),
    ('"music"  (intentionally vague)',
     "Low confidence  ·  similarity scores all cluster together",
     RGBColor(0xFF, 0xA5, 0x00)),
]

for i, (query, result, color) in enumerate(demos):
    top = 1.75 + i * 1.75
    box(s4, 0.5, top, 12.3, 1.5, fill=RGBColor(0x22, 0x22, 0x38), line=color)
    pill(s4, f"Query {i+1}", 0.65, top + 0.08, 1.3, 0.38,
         fill=color, tsize=14)
    txt(s4, query, 2.1, top + 0.05, 10.5, 0.55, size=18,
        bold=True, color=WHITE)
    txt(s4, f"→  {result}", 0.75, top + 0.75, 11.8, 0.5,
        size=17, color=color)

txt(s4, "Tip: open the RAG expander to see exactly what Claude was given",
    0.5, 7.0, 12.3, 0.4, size=15, color=LIGHT_GRAY, italic=True,
    align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — What I Learned
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5)

txt(s5, "What I Learned", 0.5, 0.25, 12, 0.65, size=32, bold=True, color=WHITE)
accent_bar(s5, top=0.95)

learnings = [
    ("🔍", "Generation is only as good as retrieval",
     "When Claude got the right songs, answers were specific and accurate.\n"
     "Weak retrieval = weak recommendation — Claude was never the bottleneck.",
     ACCENT_TEAL),
    ("📊", "The confidence score was real signal, not decoration",
     "Vague queries → clustered similarity scores → Low confidence.\n"
     "Specific queries → spread scores → High confidence. It actually worked.",
     PURPLE),
    ("⚡", "AI systems need graceful degradation",
     "Anthropic credits ran out mid-project.\n"
     "Built a Gemini 2.0 Flash fallback on the spot — app never went down.",
     RGBColor(0xFF, 0xA5, 0x00)),
]

for i, (icon, title, body, color) in enumerate(learnings):
    top = 1.15 + i * 1.9
    box(s5, 0.5, top, 12.3, 1.7, fill=RGBColor(0x22, 0x22, 0x38), line=color)
    txt(s5, icon, 0.6, top + 0.2, 0.9, 0.9, size=32, align=PP_ALIGN.CENTER)
    txt(s5, title, 1.6, top + 0.1, 10.8, 0.55, size=20, bold=True, color=color)
    txt(s5, body,  1.6, top + 0.65, 10.8, 0.9, size=16, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Portfolio & Closing
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6)

txt(s6, "Portfolio", 0.5, 0.25, 12, 0.65, size=32, bold=True, color=WHITE)
accent_bar(s6, top=0.95)

pill(s6, "github.com/yarinacs/applied-ai-system-project",
     0.5, 1.15, 7.5, 0.5, fill=RGBColor(0x25, 0x25, 0x40),
     tsize=16, tcolor=ACCENT_TEAL)

box(s6, 0.5, 1.9, 12.3, 3.8,
    fill=RGBColor(0x20, 0x20, 0x38), line=PURPLE)

quote = (
    '"This project shows that I care about building AI systems that are '
    'honest about what they\'re doing. I didn\'t just connect a language '
    'model to an input box — I built a retrieval layer that grounds every '
    'recommendation in real data, added a confidence score that tells users '
    'when the system is uncertain, and made the intermediate steps visible '
    'so anyone can see what the AI was given before it answered.\n\n'
    'Throughout this project I was drawn to the questions around transparency '
    'and reliability — not just \'does the AI answer?\' but \'does the user '
    'understand why, and can they trust it?\' That instinct is what I want '
    'to bring to every AI system I build."'
)
txt(s6, quote, 0.75, 2.1, 11.8, 3.4, size=17, color=WHITE, italic=True)

txt(s6, "Thank you  ·  Questions?", 0.5, 6.85, 12.3, 0.55,
    size=26, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "VibeMatch_AI_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
